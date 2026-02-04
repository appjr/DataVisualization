"""
Convert Class3.md markdown slides to PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re
import os

print("Creating PowerPoint presentation from Class3.md...")

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Read markdown file
with open('Class3.md', 'r', encoding='utf-8') as f:
    content = f.read()

# Split into slides
slides_content = re.split(r'\n---\n', content)

print(f"Found {len(slides_content)} slides to convert...")

def clean_text(text):
    """Remove markdown formatting"""
    # Remove bold
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    # Remove italic
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    # Remove code blocks
    text = re.sub(r'`([^`]+)`', r'\1', text)
    # Remove emojis that might cause issues
    text = text.replace('✅', '[OK]').replace('❌', '[X]').replace('⭐', '*')
    text = text.replace('✨', '*').replace('🎯', '*').replace('📊', '')
    text = text.replace('🖥️', '').replace('📱', '').replace('🎬', '')
    return text.strip()

def add_title_slide(slide, title, subtitle=''):
    """Add title and subtitle to a slide"""
    title_shape = slide.shapes.title
    title_shape.text = clean_text(title)
    
    # Format title
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    if subtitle and len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = clean_text(subtitle)
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.paragraphs[0].font.size = Pt(18)

def add_content_slide(slide, title, content_text):
    """Add title and content to a slide"""
    # Add title manually
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = clean_text(title)
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Add content
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9)
    height = Inches(5.8)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    # Parse content by lines
    lines = content_text.strip().split('\n')
    
    for line in lines:
        if not line.strip():
            continue
            
        p = text_frame.add_paragraph()
        p.text = clean_text(line)
        p.font.size = Pt(14)
        p.space_before = Pt(6)
        
        # Handle bullets
        if line.strip().startswith('-'):
            p.level = 0
            p.text = clean_text(line[1:].strip())
        elif line.strip().startswith('  -'):
            p.level = 1
            p.text = clean_text(line[3:].strip())

def add_image_to_slide(slide, image_path, title=''):
    """Add an image to a slide"""
    if title:
        # Add title textbox manually for blank layout
        left = Inches(0.5)
        top = Inches(0.2)
        width = Inches(9)
        height = Inches(0.8)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = clean_text(title)
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)
        p.alignment = PP_ALIGN.CENTER
        
        img_top = Inches(1.3)
    else:
        img_top = Inches(0.5)
    
    if os.path.exists(image_path):
        left = Inches(1)
        width = Inches(8)
        slide.shapes.add_picture(image_path, left, img_top, width=width)
    else:
        # Add placeholder text if image not found
        left = Inches(2)
        width = Inches(6)
        height = Inches(1)
        textbox = slide.shapes.add_textbox(left, img_top, width, height)
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = f"[Image: {os.path.basename(image_path)}]"
        p.alignment = PP_ALIGN.CENTER

# Process each slide
for idx, slide_content in enumerate(slides_content):
    if not slide_content.strip():
        continue
    
    # Extract title (look for ## pattern)
    title_match = re.search(r'## (Slide \d+.*?)(?:\n|$)', slide_content)
    if not title_match:
        title_match = re.search(r'## (.*?)(?:\n|$)', slide_content)
    
    if not title_match:
        # Skip non-slide content (like the header)
        continue
    
    title = title_match.group(1).strip()
    
    # Remove title from content
    content = slide_content[title_match.end():].strip()
    
    # Check if slide has an image
    image_match = re.search(r'!\[.*?\]\((.*?)\)', content)
    
    if image_match:
        image_path = image_match.group(1)
        # Remove image markdown from content
        content = re.sub(r'!\[.*?\]\(.*?\)', '', content).strip()
        
        # Create slide with image
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add image
        if not image_path.startswith('http'):
            add_image_to_slide(slide, image_path, title)
        else:
            # External link - just add title and content
            add_content_slide(slide, title, content)
    else:
        # Regular content slide
        slide_layout = prs.slide_layouts[1]  # Title and content
        slide = prs.slides.add_slide(slide_layout)
        
        # First slide special handling
        if 'Class 3 – Data Visualization' in title:
            add_title_slide(slide, title, content[:200] if len(content) > 200 else content)
        else:
            add_content_slide(slide, title, content)
    
    if (idx + 1) % 10 == 0:
        print(f"Processed {idx + 1} slides...")

# Save presentation
output_file = 'Class3_DataVisualization.pptx'
prs.save(output_file)

print(f"\n✅ PowerPoint presentation created successfully!")
print(f"📄 File: {output_file}")
print(f"📊 Total slides: {len(prs.slides)}")
print(f"\nYou can now open {output_file} in PowerPoint, Keynote, or Google Slides!")
