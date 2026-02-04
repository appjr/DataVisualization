"""
Create separate PowerPoint presentations for each major section of Class 3
"""

from pptx import Presentation
from pptx.util import Inches
import os
import shutil

print("Creating sectioned PowerPoint presentations...")

# Define major sections based on slide content
sections = {
    'Part1_Visual_Perception_Cognitive_Load': {
        'title': 'Part 1: Visual Perception & Cognitive Load',
        'slides': range(1, 22),  # Slides 1-21
        'description': 'How humans perceive visual information and process data'
    },
    'Part2_Data_Types_Encodings': {
        'title': 'Part 2: Data Types & Visual Encodings',
        'slides': range(22, 34),  # Slides 22-33
        'description': 'Understanding data types and choosing effective visual encodings'
    },
    'Part3_Grammar_of_Graphics': {
        'title': 'Part 3: Grammar of Graphics Framework',
        'slides': range(34, 45),  # Slides 34-44
        'description': 'Compositional approach to building visualizations'
    },
    'Part4_Python_Implementation': {
        'title': 'Part 4: Python Visualization Implementation',
        'slides': range(45, 85),  # Slides 45-84
        'description': 'Hands-on Python with Matplotlib and Seaborn'
    }
}

# Create section folders
print("\nCreating section folders...")
for section_name in sections.keys():
    os.makedirs(f'sections/{section_name}', exist_ok=True)

# Copy relevant slide images to each section folder
print("\nOrganizing slide images by section...")
for section_name, section_info in sections.items():
    section_dir = f'sections/{section_name}'
    
    # Copy slide images
    for slide_num in section_info['slides']:
        src_file = f'slide_images/slide_{slide_num:03d}.png'
        if os.path.exists(src_file):
            dest_file = f'{section_dir}/slide_{slide_num:03d}.png'
            shutil.copy2(src_file, dest_file)
    
    print(f"  {section_name}: {len(section_info['slides'])} slides")

# Create PowerPoint for each section
print("\nCreating PowerPoint presentations for each section...")

for section_name, section_info in sections.items():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Add title slide
    slide_layout = prs.slide_layouts[0]  # Title slide
    title_slide = prs.slides.add_slide(slide_layout)
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = section_info['title']
    subtitle.text = section_info['description'] + f"\n\n{len(section_info['slides'])} slides"
    
    # Add content slides
    for slide_num in section_info['slides']:
        img_path = f'slide_images/slide_{slide_num:03d}.png'
        
        if os.path.exists(img_path):
            # Use blank layout
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            
            # Add image filling the entire slide
            left = Inches(0)
            top = Inches(0)
            width = Inches(10)
            height = Inches(7.5)
            
            slide.shapes.add_picture(img_path, left, top, width=width, height=height)
    
    # Save section PowerPoint
    output_file = f'sections/{section_name}/{section_name}.pptx'
    prs.save(output_file)
    
    print(f"  ✅ Created: {section_name}.pptx ({len(section_info['slides'])} slides)")

# Create a master index/navigation PowerPoint
print("\nCreating Master Navigation PowerPoint...")

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Title slide
slide_layout = prs.slide_layouts[0]
title_slide = prs.slides.add_slide(slide_layout)
title = title_slide.shapes.title
subtitle = title_slide.placeholders[1]

title.text = "Class 3: Data Visualization"
subtitle.text = "Visual Perception, Cognitive Load & Python Fundamentals\n\nClick to navigate to each section"

# Create navigation slide
slide_layout = prs.slide_layouts[1]
nav_slide = prs.slides.add_slide(slide_layout)
nav_title = nav_slide.shapes.title
nav_title.text = "Course Sections - Navigation"

# Add text box with section descriptions
left = Inches(1)
top = Inches(2)
width = Inches(8)
height = Inches(4.5)

textbox = nav_slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.word_wrap = True

for idx, (section_name, section_info) in enumerate(sections.items(), 1):
    p = text_frame.add_paragraph()
    p.text = f"{idx}. {section_info['title']}"
    p.font.size = Inches(0.25)
    p.font.bold = True
    p.space_after = Inches(0.05)
    
    p = text_frame.add_paragraph()
    p.text = f"   {section_info['description']}"
    p.font.size = Inches(0.18)
    p.space_after = Inches(0.15)
    
    p = text_frame.add_paragraph()
    p.text = f"   📊 {len(section_info['slides'])} slides | File: {section_name}.pptx"
    p.font.size = Inches(0.15)
    p.font.italic = True
    p.space_after = Inches(0.25)

# Save master navigation
master_file = 'sections/00_Master_Navigation.pptx'
prs.save(master_file)

print(f"\n✅ Created Master Navigation: 00_Master_Navigation.pptx")

# Create a README for the sections
readme_content = f"""# Class 3 - Sectioned Presentations

This directory contains the Class 3 content divided into 4 major sections for easier navigation and teaching.

## Master Navigation
- **00_Master_Navigation.pptx** - Overview and links to all sections

## Section Breakdown

### Part 1: Visual Perception & Cognitive Load (Slides 1-21)
**File**: `Part1_Visual_Perception_Cognitive_Load/Part1_Visual_Perception_Cognitive_Load.pptx`
**Topics**:
- Preattentive processing and the 200ms window
- Visual attention limitations
- Change blindness and inattentional blindness
- Gestalt principles of grouping
- Cognitive load theory (intrinsic, extraneous, germane)
- Dashboard design principles

**Key Learning**: Understand how humans perceive and process visual information

---

### Part 2: Data Types & Visual Encodings (Slides 22-33)
**File**: `Part2_Data_Types_Encodings/Part2_Data_Types_Encodings.pptx`
**Topics**:
- Four fundamental data types (Nominal, Ordinal, Quantitative, Temporal)
- Bertin's visual variables
- Cleveland & McGill's ranking of visual encodings
- Why color is weak for quantitative data
- Data × Task × Encoding framework
- Common encoding violations

**Key Learning**: Match data types to appropriate visual encodings

---

### Part 3: Grammar of Graphics Framework (Slides 34-44)
**File**: `Part3_Grammar_of_Graphics/Part3_Grammar_of_Graphics.pptx`
**Topics**:
- Introduction to Grammar of Graphics
- Seven core components (Data, Aesthetics, Geoms, Stats, Scales, Coordinates, Facets)
- Compositional thinking vs. chart templates
- How grammar reduces cognitive load
- Systematic approach to visualization design

**Key Learning**: Think in layers and mappings, not chart types

---

### Part 4: Python Visualization Implementation (Slides 45-84)
**File**: `Part4_Python_Implementation/Part4_Python_Implementation.pptx`
**Topics**:
- Python visualization ecosystem
- Matplotlib anatomy and fundamentals
- Seaborn for statistical graphics
- Applying perception principles in code
- Color palette selection
- Reducing clutter and cognitive load
- Direct labeling vs. legends
- Complete examples and exercises

**Key Learning**: Implement perception-based design in Python

---

## Usage

### For Teaching Full Class:
Use the complete `Class3_ImageSlides.pptx` (84 slides)

### For Modular Teaching:
1. Start with `00_Master_Navigation.pptx` to show course structure
2. Teach each part separately over multiple sessions
3. Each section is self-contained with 10-40 slides

### For Self-Study:
Work through parts sequentially:
1. Visual Perception (foundational concepts)
2. Data Types & Encodings (theory)
3. Grammar of Graphics (framework)
4. Python Implementation (hands-on practice)

## File Structure

```
sections/
├── 00_Master_Navigation.pptx
├── Part1_Visual_Perception_Cognitive_Load/
│   ├── Part1_Visual_Perception_Cognitive_Load.pptx
│   └── slide_001.png through slide_021.png
├── Part2_Data_Types_Encodings/
│   ├── Part2_Data_Types_Encodings.pptx
│   └── slide_022.png through slide_033.png
├── Part3_Grammar_of_Graphics/
│   ├── Part3_Grammar_of_Graphics.pptx
│   └── slide_034.png through slide_044.png
└── Part4_Python_Implementation/
    ├── Part4_Python_Implementation.pptx
    └── slide_045.png through slide_084.png
```

## Benefits of Sectioned Approach

✅ **Modular Teaching** - Teach over multiple sessions  
✅ **Focused Review** - Students can review specific topics  
✅ **Easier Navigation** - Jump to relevant content quickly  
✅ **Flexible Scheduling** - Adapt to different class lengths  
✅ **Self-Paced Learning** - Students work through sections independently  

---

**Total**: 84 slides across 4 sections + 1 master navigation
"""

with open('sections/README.md', 'w') as f:
    f.write(readme_content)

print("✅ Created sections/README.md\n")

# Summary
print("="*60)
print("SECTIONED PRESENTATIONS CREATED SUCCESSFULLY!")
print("="*60)
print(f"\nCreated {len(sections)} section PowerPoints:")
for section_name, section_info in sections.items():
    print(f"  • {section_info['title']}: {len(section_info['slides'])} slides")

print(f"\n📁 All files saved in: sections/")
print(f"📄 Master navigation: sections/00_Master_Navigation.pptx")
print(f"📖 Documentation: sections/README.md")
print(f"\nEach section folder contains:")
print(f"  - PowerPoint presentation")
print(f"  - Individual slide images")
