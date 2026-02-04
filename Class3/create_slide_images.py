"""
Render each markdown slide as an image, then create PowerPoint with one image per slide
"""

import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.patches import FancyBboxPatch
from PIL import Image
import re
import os
from pptx import Presentation
from pptx.util import Inches

print("Creating slide images from Class3.md...")

# Create output directory
os.makedirs('slide_images', exist_ok=True)

# Read markdown file
with open('Class3.md', 'r', encoding='utf-8') as f:
    content = f.read()

# Split into slides
slides_content = re.split(r'\n---\n', content)

print(f"Found {len(slides_content)} slides to convert...")

def clean_text(text):
    """Remove markdown formatting for display"""
    # Remove bold
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    # Remove italic
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    # Remove code backticks
    text = re.sub(r'`([^`]+)`', r'\1', text)
    # Escape dollar signs to prevent mathtext interpretation
    text = text.replace('$', r'\$')
    return text.strip()

def wrap_text(text, max_chars=80):
    """Wrap text to fit within slide"""
    words = text.split()
    lines = []
    current_line = []
    current_length = 0
    
    for word in words:
        if current_length + len(word) + 1 <= max_chars:
            current_line.append(word)
            current_length += len(word) + 1
        else:
            if current_line:
                lines.append(' '.join(current_line))
            current_line = [word]
            current_length = len(word)
    
    if current_line:
        lines.append(' '.join(current_line))
    
    return lines

def render_slide_to_image(slide_num, title, content, has_image=False, image_path=None):
    """Render a slide to an image file"""
    fig = plt.figure(figsize=(10, 7.5), facecolor='white')
    ax = fig.add_subplot(111)
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 7.5)
    ax.axis('off')
    
    # Add title
    ax.text(5, 7, clean_text(title), fontsize=20, fontweight='bold',
           ha='center', va='top', color='#003366', wrap=True)
    
    # Add content or image
    if has_image and image_path and os.path.exists(image_path):
        # Load and display image
        try:
            img = Image.open(image_path)
            # Position image below title
            ax.imshow(img, extent=[1, 9, 0.5, 6], aspect='auto')
        except Exception as e:
            print(f"  Warning: Could not load image {image_path}: {e}")
            # Fall back to text content
            add_text_content(ax, content)
    else:
        add_text_content(ax, content)
    
    # Save slide image
    output_file = f'slide_images/slide_{slide_num:03d}.png'
    plt.savefig(output_file, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    
    return output_file

def add_text_content(ax, content):
    """Add text content to slide"""
    # Clean and parse content
    lines = content.strip().split('\n')
    
    y_pos = 6.2
    line_height = 0.25
    
    for line in lines[:25]:  # Limit to 25 lines to fit on slide
        if not line.strip():
            y_pos -= line_height * 0.5
            continue
        
        cleaned = clean_text(line)
        
        # Check for bullet points
        if line.strip().startswith('-'):
            # Bullet point
            bullet_text = '• ' + cleaned[1:].strip()
            wrapped = wrap_text(bullet_text, 70)
            for wrapped_line in wrapped:
                ax.text(0.8, y_pos, wrapped_line, fontsize=11,
                       ha='left', va='top', wrap=True)
                y_pos -= line_height
        elif line.strip().startswith('  -'):
            # Sub-bullet
            bullet_text = '  ◦ ' + cleaned[2:].strip()
            wrapped = wrap_text(bullet_text, 65)
            for wrapped_line in wrapped:
                ax.text(1.2, y_pos, wrapped_line, fontsize=10,
                       ha='left', va='top', wrap=True)
                y_pos -= line_height
        elif line.strip().startswith('#'):
            # Subheading
            wrapped = wrap_text(cleaned.replace('#', '').strip(), 70)
            for wrapped_line in wrapped:
                ax.text(5, y_pos, wrapped_line, fontsize=13, fontweight='bold',
                       ha='center', va='top', color='#003366', wrap=True)
                y_pos -= line_height * 1.2
        else:
            # Regular text
            wrapped = wrap_text(cleaned, 75)
            for wrapped_line in wrapped:
                ax.text(0.8, y_pos, wrapped_line, fontsize=11,
                       ha='left', va='top', wrap=True)
                y_pos -= line_height
        
        if y_pos < 0.5:
            break

# Process each slide
slide_images = []
slide_count = 0

for idx, slide_content in enumerate(slides_content):
    if not slide_content.strip():
        continue
    
    # Extract title
    title_match = re.search(r'## (Slide \d+.*?)(?:\n|$)', slide_content)
    if not title_match:
        title_match = re.search(r'## (.*?)(?:\n|$)', slide_content)
    
    if not title_match:
        continue
    
    title = title_match.group(1).strip()
    content = slide_content[title_match.end():].strip()
    
    # Check for image
    image_match = re.search(r'!\[.*?\]\((.*?)\)', content)
    has_image = bool(image_match)
    image_path = image_match.group(1) if image_match else None
    
    # Remove image markdown from content
    if has_image:
        content = re.sub(r'!\[.*?\]\(.*?\)', '', content).strip()
    
    # Render slide
    slide_count += 1
    output_file = render_slide_to_image(slide_count, title, content, has_image, image_path)
    slide_images.append(output_file)
    
    if slide_count % 10 == 0:
        print(f"Rendered {slide_count} slides...")

print(f"\n✅ Rendered {slide_count} slides to images!")
print(f"📁 Images saved in: slide_images/")

# Now create PowerPoint with one image per slide
print("\nCreating PowerPoint presentation from slide images...")

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

for idx, img_path in enumerate(slide_images):
    # Use blank layout
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Add image filling the entire slide
    left = Inches(0)
    top = Inches(0)
    width = Inches(10)
    height = Inches(7.5)
    
    slide.shapes.add_picture(img_path, left, top, width=width, height=height)
    
    if (idx + 1) % 20 == 0:
        print(f"Added {idx + 1} slides to PowerPoint...")

# Save presentation
output_pptx = 'Class3_ImageSlides.pptx'
prs.save(output_pptx)

print(f"\n✅ PowerPoint with image slides created successfully!")
print(f"📄 File: {output_pptx}")
print(f"📊 Total slides: {len(slide_images)}")
print(f"🖼️  Slide images: slide_images/ directory")
print(f"\nYou now have:")
print(f"  1. Individual slide images in 'slide_images/' folder")
print(f"  2. PowerPoint presentation '{output_pptx}' with one image per slide")
