"""
Create Jetsons-inspired slide images for Class 4 markdown slides,
then build a PowerPoint with one image per slide.
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import List, Tuple

import numpy as np
import matplotlib.pyplot as plt
from matplotlib.patches import Circle, Rectangle, FancyBboxPatch, Polygon
from pptx import Presentation
from pptx.util import Inches


BASE_DIR = Path(__file__).resolve().parent
MD_FILE = BASE_DIR / "Class4.md"
OUTPUT_DIR = BASE_DIR / "slide_images"
OUTPUT_PPTX = BASE_DIR / "Class4_Jetsons_ImageSlides.pptx"


def split_markdown_slides(content: str) -> List[str]:
    """Split markdown by slide separators."""
    parts = re.split(r"\n---\n", content)
    return [p.strip() for p in parts if p.strip()]


def extract_title_and_points(slide_text: str) -> Tuple[str, List[str]]:
    """Extract a title plus a few bullet points from a markdown slide."""
    lines = [ln.strip() for ln in slide_text.splitlines() if ln.strip()]

    title = "Class 4 Slide"
    h2 = next((ln[3:].strip() for ln in lines if ln.startswith("## ")), None)
    h1 = next((ln[2:].strip() for ln in lines if ln.startswith("# ")), None)
    if h2:
        title = h2
    elif h1:
        title = h1

    bullets: List[str] = []
    for ln in lines:
        if ln.startswith("- "):
            clean = re.sub(r"\*\*(.*?)\*\*", r"\1", ln[2:]).strip()
            clean = re.sub(r"`([^`]+)`", r"\1", clean)
            bullets.append(clean)

    if not bullets:
        # fallback from plain lines
        for ln in lines:
            if ln.startswith("#") or ln.startswith("!") or ln.startswith("```"):
                continue
            if len(ln) > 20:
                bullets.append(re.sub(r"[*`]+", "", ln))
            if len(bullets) >= 3:
                break

    return title, bullets[:3]


def detect_theme(title: str, bullets: List[str], slide_text: str) -> str:
    """Classify slide into a visual theme based on content keywords."""
    all_text = " ".join([title] + bullets + [slide_text]).lower()
    title_text = title.lower()

    checks = {
        "missing": ["missing", "mcar", "mar", "mnar", "imputation", "null"],
        "correlation": ["correlation", "heatmap", "pair plot", "anscombe", "multivariate"],
        "distribution": ["distribution", "histogram", "kde", "skew", "univariate"],
        "outlier": ["outlier", "iqr", "z-score", "percentile"],
        "comparison": ["categorical", "group", "comparison", "box plot", "violin", "scatter"],
        "workflow": ["workflow", "lifecycle", "process", "checklist", "template", "pipeline"],
        "exercise": ["exercise", "your task", "deliverable", "time:"],
        "code": ["python", "code", "import", "def ", "pd.", "sns.", "plt."],
        "summary": ["summary", "takeaways", "best practices", "mistakes", "tips", "final thoughts"],
        "assignment": ["assignment", "due", "grading", "rubric", "next class", "resources"],
    }

    scores = {k: 0 for k in checks}
    for theme, kws in checks.items():
        for kw in kws:
            if kw in all_text:
                scores[theme] += 1
            if kw in title_text:
                scores[theme] += 2

    # Prefer code/exercise/assignment when explicit in title
    for strong in ["exercise", "code", "assignment", "summary"]:
        if scores[strong] >= 2:
            return strong

    best_theme, best_score = "general", 0
    for k, v in scores.items():
        if v > best_score:
            best_theme, best_score = k, v
    return best_theme if best_score > 0 else "general"


def _draw_panel_frame(ax) -> None:
    panel = FancyBboxPatch(
        (0.55, 1.25),
        8.9,
        4.05,
        boxstyle="round,pad=0.04,rounding_size=0.14",
        linewidth=1.4,
        edgecolor="#9af2ff",
        facecolor="#0e1e48",
        alpha=0.84,
        zorder=9,
    )
    ax.add_patch(panel)


def _draw_bar_chart(ax, x0, y0, w, h, rng, color="#8fe8ff"):
    bars = rng.integers(5, 9)
    bw = w / (bars + 1)
    for i in range(bars):
        bh = float(rng.uniform(0.2, 0.95) * h)
        ax.add_patch(Rectangle((x0 + i * bw + 0.05, y0), bw * 0.65, bh, color=color, alpha=0.8, zorder=12))


def _draw_line_chart(ax, x0, y0, w, h, rng, color="#ffd166"):
    xs = np.linspace(x0 + 0.05, x0 + w - 0.05, 14)
    vals = rng.uniform(0.15, 0.92, size=len(xs))
    ys = y0 + vals * h
    ax.plot(xs, ys, color=color, linewidth=2.0, zorder=12)
    ax.scatter(xs, ys, s=10, color="#efffff", zorder=13)


def _draw_heatmap(ax, x0, y0, w, h, rows, cols, rng):
    cw = w / cols
    ch = h / rows
    for r in range(rows):
        for c in range(cols):
            val = rng.uniform()
            col = (0.25 + 0.65 * val, 0.35, 0.85 - 0.5 * val)
            ax.add_patch(Rectangle((x0 + c * cw, y0 + r * ch), cw - 0.01, ch - 0.01, color=col, alpha=0.9, zorder=12))


def _draw_theme_content(ax, theme: str, bullets: List[str], rng: np.random.Generator) -> None:
    _draw_panel_frame(ax)

    # Header for content area
    ax.text(0.9, 4.95, f"Content Theme: {theme.title()}", fontsize=10.5, color="#9ef8ff", zorder=13, fontweight="bold")

    if theme == "distribution":
        _draw_bar_chart(ax, 0.95, 1.75, 3.8, 2.2, rng, color="#79dbff")
        _draw_line_chart(ax, 5.2, 1.75, 3.8, 2.2, rng, color="#ffb866")
        ax.text(1.0, 4.5, "Histograms / KDE / Skewness", color="#e5fcff", fontsize=10, zorder=13)

    elif theme == "correlation":
        _draw_heatmap(ax, 1.0, 1.7, 3.9, 2.5, 6, 6, rng)
        _draw_line_chart(ax, 5.25, 1.8, 3.7, 1.0, rng, color="#a8ff8b")
        _draw_bar_chart(ax, 5.25, 3.1, 3.7, 1.1, rng, color="#bfa2ff")
        ax.text(1.0, 4.5, "Correlation Matrix / Pair Relationships", color="#e5fcff", fontsize=10, zorder=13)

    elif theme == "missing":
        _draw_heatmap(ax, 1.0, 1.7, 4.2, 2.5, 8, 10, rng)
        for i, t in enumerate(["MCAR", "MAR", "MNAR"]):
            ax.add_patch(FancyBboxPatch((5.6, 3.8 - i * 0.8), 3.3, 0.55, boxstyle="round,pad=0.02,rounding_size=0.08",
                                        facecolor="#273777", edgecolor="#9cf8ff", linewidth=1.0, zorder=12, alpha=0.9))
            ax.text(5.75, 4.08 - i * 0.8, t, fontsize=10, color="#f1ffff", zorder=13, fontweight="bold")
        ax.text(1.0, 4.5, "Missingness Patterns / Imputation", color="#e5fcff", fontsize=10, zorder=13)

    elif theme == "outlier":
        _draw_bar_chart(ax, 0.95, 1.75, 4.0, 2.2, rng, color="#9be7ff")
        xs = np.linspace(5.4, 9.0, 65)
        ys = 2.8 + 0.45 * np.sin(xs * 1.5) + rng.normal(0, 0.07, len(xs))
        ax.scatter(xs, ys, s=10, color="#d9f7ff", alpha=0.8, zorder=12)
        ax.scatter([8.6, 8.8], [3.9, 4.2], s=24, color="#ff6b8f", zorder=13)
        ax.text(5.3, 4.5, "IQR / Z-score Outlier Flags", color="#e5fcff", fontsize=10, zorder=13)

    elif theme == "comparison":
        _draw_bar_chart(ax, 1.0, 1.75, 3.8, 2.2, rng, color="#8be7af")
        _draw_bar_chart(ax, 5.2, 1.75, 3.8, 2.2, rng, color="#87d8ff")
        ax.text(1.0, 4.5, "Group Comparison / Categorical Analysis", color="#e5fcff", fontsize=10, zorder=13)

    elif theme == "workflow":
        nodes = [
            (1.0, 3.7, "Data"),
            (2.8, 3.7, "Quality"),
            (4.6, 3.7, "EDA"),
            (6.4, 3.7, "Features"),
            (8.2, 3.7, "Model"),
        ]
        for i, (x, y, label) in enumerate(nodes):
            ax.add_patch(FancyBboxPatch((x, y), 1.2, 0.6, boxstyle="round,pad=0.02,rounding_size=0.08",
                                        facecolor="#2e4f86", edgecolor="#a8fbff", linewidth=1.1, zorder=12, alpha=0.9))
            ax.text(x + 0.6, y + 0.3, label, ha="center", va="center", fontsize=9.2, color="#f4ffff", zorder=13)
            if i < len(nodes) - 1:
                ax.arrow(x + 1.25, y + 0.3, 0.45, 0, head_width=0.08, head_length=0.1,
                         fc="#b0fbff", ec="#b0fbff", zorder=13, linewidth=1.0)
        ax.text(1.0, 2.3, "Question-driven iterative workflow", color="#e5fcff", fontsize=10, zorder=13)

    elif theme == "exercise":
        for i in range(4):
            ax.add_patch(FancyBboxPatch((1.0 + i * 2.05, 2.1), 1.8, 1.8, boxstyle="round,pad=0.03,rounding_size=0.1",
                                        facecolor="#233e79", edgecolor="#9cf8ff", linewidth=1.0, zorder=12, alpha=0.9))
            _draw_bar_chart(ax, 1.1 + i * 2.05, 2.25, 1.5, 1.3, rng, color="#ffc178")
        ax.text(1.0, 4.5, "Hands-on mini analysis challenges", color="#e5fcff", fontsize=10, zorder=13)

    elif theme == "code":
        ax.add_patch(Rectangle((0.95, 1.75), 8.0, 2.45, color="#09162f", alpha=0.95, zorder=12))
        code_lines = ["import pandas as pd", "import seaborn as sns", "corr = df.corr()", "sns.heatmap(corr)", "plt.show()"]
        for i, line in enumerate(code_lines):
            ax.text(1.15, 3.95 - i * 0.45, line, fontsize=9.3, color="#87f7c5", family="monospace", zorder=13)
        ax.text(1.0, 4.5, "Python EDA implementation", color="#e5fcff", fontsize=10, zorder=13)

    elif theme in {"summary", "assignment"}:
        for i in range(5):
            y = 4.1 - i * 0.5
            ax.add_patch(Circle((1.2, y), 0.05, color="#8af6ff", zorder=13))
            ax.plot([1.35, 8.9], [y, y], color="#7fe1ff", alpha=0.65, linewidth=1.0, zorder=12)
        ax.text(1.0, 4.5, "Key points / deliverables", color="#e5fcff", fontsize=10, zorder=13)

    else:
        _draw_line_chart(ax, 1.0, 1.75, 4.0, 2.2, rng, color="#95e5ff")
        _draw_bar_chart(ax, 5.1, 1.75, 3.8, 2.2, rng, color="#ffc578")
        ax.text(1.0, 4.5, "General EDA Concepts", color="#e5fcff", fontsize=10, zorder=13)

    # content snippets (ensures each slide carries specific text)
    y = 1.55
    for b in bullets[:2]:
        txt = re.sub(r"\s+", " ", b).strip()
        if len(txt) > 85:
            txt = txt[:82] + "..."
        txt = txt.replace("$", "\\$")
        ax.text(0.95, y, f"• {txt}", fontsize=9.3, color="#e7feff", zorder=13)
        y -= 0.28


def _gradient_background(ax, rng: np.random.Generator) -> None:
    h, w = 600, 800
    y = np.linspace(0, 1, h).reshape(-1, 1)
    top = np.array([0.08, 0.12, 0.32])
    bottom = np.array([0.50, 0.30, 0.62])
    grad = top * (1 - y) + bottom * y
    grad = np.repeat(grad[:, None, :], w, axis=1)
    ax.imshow(grad, extent=[0, 10, 0, 7.5], aspect="auto", zorder=0)

    # stars
    xs = rng.uniform(0.2, 9.8, 120)
    ys = rng.uniform(3.5, 7.3, 120)
    sizes = rng.uniform(5, 22, 120)
    ax.scatter(xs, ys, s=sizes, c="#f4f8ff", alpha=0.8, linewidths=0, zorder=1)


def _retro_city(ax, rng: np.random.Generator) -> None:
    # ground glow
    ax.add_patch(Rectangle((0, 0), 10, 1.1, color="#1b1737", alpha=0.95, zorder=2))
    ax.add_patch(Rectangle((0, 1.0), 10, 0.15, color="#73d6ff", alpha=0.45, zorder=3))

    # buildings + domes
    x = 0.3
    while x < 9.7:
        bw = rng.uniform(0.4, 0.9)
        bh = rng.uniform(0.8, 2.3)
        color = rng.choice(["#2d2f64", "#413a84", "#30407c", "#275d86"])
        ax.add_patch(Rectangle((x, 1.1), bw, bh, color=color, alpha=0.92, zorder=4))

        dome_r = bw * rng.uniform(0.25, 0.45)
        ax.add_patch(Circle((x + bw / 2, 1.1 + bh), dome_r, color="#7fe3ff", alpha=0.55, zorder=5))

        # windows
        n_win = int(max(2, bh / 0.35))
        for i in range(n_win):
            wy = 1.25 + i * 0.32
            if wy < 1.1 + bh - 0.15:
                ax.add_patch(Rectangle((x + 0.08, wy), bw - 0.16, 0.05, color="#a4f7ff", alpha=0.6, zorder=6))

        x += bw + rng.uniform(0.12, 0.25)

    # monorail arc
    rail_x = np.linspace(0, 10, 200)
    rail_y = 2.7 + 0.35 * np.sin(rail_x * 0.7)
    ax.plot(rail_x, rail_y, color="#88f2ff", linewidth=3.0, alpha=0.8, zorder=6)


def _flying_car(ax, x: float, y: float, scale: float, color: str) -> None:
    body = FancyBboxPatch(
        (x, y),
        0.9 * scale,
        0.28 * scale,
        boxstyle="round,pad=0.02,rounding_size=0.12",
        linewidth=1.2,
        edgecolor="#f7ffff",
        facecolor=color,
        alpha=0.95,
        zorder=8,
    )
    ax.add_patch(body)
    canopy = Polygon(
        [
            (x + 0.15 * scale, y + 0.28 * scale),
            (x + 0.55 * scale, y + 0.48 * scale),
            (x + 0.75 * scale, y + 0.28 * scale),
        ],
        closed=True,
        facecolor="#d2f8ff",
        edgecolor="#f7ffff",
        linewidth=1.0,
        alpha=0.9,
        zorder=9,
    )
    ax.add_patch(canopy)
    ax.plot([x - 0.55 * scale, x], [y + 0.12 * scale, y + 0.12 * scale], color="#8cf7ff", linewidth=2, alpha=0.7, zorder=7)


def create_jetsons_slide_image(
    slide_num: int,
    title: str,
    bullets: List[str],
    slide_text: str,
    out_path: Path,
) -> None:
    rng = np.random.default_rng(slide_num * 7919 + len(title) * 101)
    fig, ax = plt.subplots(figsize=(10, 7.5), dpi=150)
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 7.5)
    ax.axis("off")

    _gradient_background(ax, rng)

    # planet/moon for retro-futuristic mood
    moon_x = rng.uniform(7.4, 9.2)
    moon_y = rng.uniform(5.7, 6.8)
    moon_r = rng.uniform(0.35, 0.55)
    ax.add_patch(Circle((moon_x, moon_y), moon_r, color="#ffe39a", alpha=0.75, zorder=2))

    _retro_city(ax, rng)

    theme = detect_theme(title, bullets, slide_text)

    # flying cars
    car_colors = ["#ff7ea9", "#6fe3ff", "#ffd166", "#b18cff"]
    for _ in range(3):
        _flying_car(
            ax,
            x=float(rng.uniform(0.8, 8.3)),
            y=float(rng.uniform(3.0, 5.2)),
            scale=float(rng.uniform(0.8, 1.2)),
            color=str(rng.choice(car_colors)),
        )

    # Title panel
    panel = FancyBboxPatch(
        (0.55, 5.5),
        8.9,
        1.35,
        boxstyle="round,pad=0.05,rounding_size=0.12",
        linewidth=1.8,
        edgecolor="#98f6ff",
        facecolor="#11183d",
        alpha=0.82,
        zorder=10,
    )
    ax.add_patch(panel)

    safe_title = title.replace("$", "\\$")
    if len(safe_title) > 80:
        safe_title = safe_title[:77] + "..."

    ax.text(
        0.85,
        6.28,
        f"Slide {slide_num:02d}",
        fontsize=10,
        color="#a6fbff",
        fontweight="bold",
        va="center",
        zorder=11,
    )
    ax.text(
        0.85,
        5.93,
        safe_title,
        fontsize=20,
        color="#f4feff",
        fontweight="bold",
        va="center",
        zorder=11,
    )

    # content-specific visual composition
    _draw_theme_content(ax, theme, bullets, rng)

    # keyword chips from bullets
    chip_y = 5.0
    x = 0.8
    for b in bullets[:3]:
        chip_text = re.sub(r"\s+", " ", b).strip()
        if len(chip_text) > 44:
            chip_text = chip_text[:41] + "..."
        w = min(3.6, 0.42 + 0.07 * len(chip_text))
        chip = FancyBboxPatch(
            (x, chip_y),
            w,
            0.42,
            boxstyle="round,pad=0.02,rounding_size=0.10",
            linewidth=1.0,
            edgecolor="#9cf5ff",
            facecolor="#1f2f62",
            alpha=0.78,
            zorder=10,
        )
        ax.add_patch(chip)
        ax.text(x + 0.11, chip_y + 0.21, chip_text, fontsize=8.8, color="#dffcff", va="center", zorder=11)
        x += w + 0.16
        if x > 9.0:
            break

    plt.tight_layout(pad=0)
    fig.savefig(out_path, facecolor="white", bbox_inches="tight")
    plt.close(fig)


def build_pptx_from_images(image_paths: List[Path], output_pptx: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for img_path in image_paths:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(img_path), Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))

    prs.save(str(output_pptx))


def main() -> None:
    if not MD_FILE.exists():
        raise FileNotFoundError(f"Could not find markdown file: {MD_FILE}")

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    content = MD_FILE.read_text(encoding="utf-8")
    slides = split_markdown_slides(content)
    image_paths: List[Path] = []

    print(f"Found {len(slides)} markdown slides in {MD_FILE.name}")
    print("Generating Jetsons-inspired slide images...")

    for i, slide_text in enumerate(slides, start=1):
        title, bullets = extract_title_and_points(slide_text)
        out_path = OUTPUT_DIR / f"slide_{i:03d}.png"
        create_jetsons_slide_image(i, title, bullets, slide_text, out_path)
        image_paths.append(out_path)

        if i % 10 == 0 or i == len(slides):
            print(f"  Created {i}/{len(slides)} images")

    print(f"Creating PowerPoint: {OUTPUT_PPTX.name}")
    build_pptx_from_images(image_paths, OUTPUT_PPTX)

    print("\nDone!")
    print(f"  Images: {OUTPUT_DIR}")
    print(f"  PPTX:   {OUTPUT_PPTX}")
    print(f"  Total slides/images: {len(image_paths)}")


if __name__ == "__main__":
    main()
