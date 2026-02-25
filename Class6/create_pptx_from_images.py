#!/usr/bin/env python3
"""
Create a PowerPoint presentation from slide images.
Each image is placed on its own slide.
"""

from pptx import Presentation
from pptx.util import Inches
from pathlib import Path
import os

def create_pptx_from_images(image_dir, output_file):
    """
    Create a PPTX file with one image per slide.
    
    Args:
        image_dir: Directory containing the slide images
        output_file: Output PPTX filename
    """
    # Create a presentation object
    prs = Presentation()
    
    # Set slide dimensions (16:9 widescreen)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Get all PNG files in the directory, sorted
    image_dir_path = Path(image_dir)
    image_files = sorted([f for f in image_dir_path.glob('slide_*_final.png')])
    
    print(f"Found {len(image_files)} slide images")
    
    # Add each image as a new slide
    for i, image_file in enumerate(image_files, 1):
        print(f"Adding slide {i}/{len(image_files)}: {image_file.name}")
        
        # Add a blank slide
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add the image to fill the entire slide
        left = Inches(0)
        top = Inches(0)
        width = prs.slide_width
        height = prs.slide_height
        
        slide.shapes.add_picture(
            str(image_file),
            left, top,
            width=width,
            height=height
        )
    
    # Save the presentation
    prs.save(output_file)
    print(f"\nPresentation created successfully: {output_file}")
    print(f"Total slides: {len(image_files)}")

if __name__ == "__main__":
    # Define paths
    script_dir = Path(__file__).parent
    image_dir = script_dir / "slide_images"
    output_file = script_dir / "Class6_Slides.pptx"
    
    # Create the presentation
    create_pptx_from_images(image_dir, output_file)
