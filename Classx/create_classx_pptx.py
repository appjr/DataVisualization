#!/usr/bin/env python3
"""
Create Classx_Slides.pptx from the generated slide images.
Each image is placed on its own slide (one image per slide, 80 total).
"""

from pptx import Presentation
from pptx.util import Inches
from pathlib import Path


def create_pptx_from_images(image_dir, output_file):
    prs = Presentation()

    # 16:9 widescreen dimensions
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    image_dir_path = Path(image_dir)
    image_files = sorted(image_dir_path.glob('slide_*_final.png'))

    print(f"Found {len(image_files)} slide images")

    blank_layout = prs.slide_layouts[6]  # completely blank

    for i, img_path in enumerate(image_files, 1):
        print(f"  Adding slide {i:>2}/{len(image_files)}: {img_path.name}")
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path),
            left=Inches(0), top=Inches(0),
            width=prs.slide_width, height=prs.slide_height
        )

    prs.save(output_file)
    print(f"\n✅ Saved: {output_file}  ({len(image_files)} slides)")


if __name__ == "__main__":
    script_dir  = Path(__file__).parent
    image_dir   = script_dir / "slide_images"
    output_file = script_dir / "Classx_Slides.pptx"
    create_pptx_from_images(image_dir, output_file)
